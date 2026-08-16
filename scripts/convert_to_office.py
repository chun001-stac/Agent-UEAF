#!/usr/bin/env python3
"""Convert a portion of the synthetic TXT corpus into Office formats.

Keeps the TXT source files and adds docx / pdf / pptx / xlsx variants per a
per-category format mix (docx > pdf > pptx/xlsx, roughly matching the manifest
enterprise format ratios). Only synthetic docs (license == synthetic-corpus)
are converted; real PDFs/RFCs stay untouched.

Deterministic: same seed + (category, seq) -> same target format, and a given
target file is only written once (idempotent).

Formats:
  - .docx  via python-docx
  - .pptx  via python-pptx
  - .xlsx  via openpyxl
  - .pdf   via reportlab (CID font STSong-Light for Chinese)

Usage:
  python3 scripts/convert_to_office.py                 # convert all categories
  python3 scripts/convert_to_office.py --category 03-policy-sop
  python3 scripts/convert_to_office.py --dry-run
"""

from __future__ import annotations

import argparse
import html
import json
import random
import re
import zlib
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SOURCES_JSON = ROOT / "document" / "sources.json"
DEFAULT_RAW = ROOT / "document" / "raw"
DEFAULT_SEED = 20260816

PREFIX = {
    "01-product": "product", "02-tech-api": "tech", "03-policy-sop": "policy",
    "04-sales-pre": "sales", "05-faq-support": "faq", "06-project-delivery": "project",
    "07-training-hr": "training", "08-meeting-decision": "meeting",
    "09-version-conflict": "version", "10-noise": "noise",
}

# Semantic relevance of each format for each category (0 = not applicable,
# higher = more natural). The generator's IPF allocation uses these weights to
# hit the manifest format targets exactly while keeping format-to-category
# mapping sensible (policies->docx, sales/training->pptx, FAQ->xlsx/csv, ...).
RELEVANCE: dict[str, dict[str, float]] = {
    "01-product": {"docx": 4, "pdf": 4, "pptx": 2, "xlsx": 2, "html": 3, "md": 2, "json": 0, "csv": 0, "txt": 1},
    "02-tech-api": {"docx": 2, "pdf": 4, "pptx": 1, "xlsx": 2, "html": 4, "md": 5, "json": 4, "csv": 1, "txt": 1},
    "03-policy-sop": {"docx": 5, "pdf": 3, "pptx": 1, "xlsx": 2, "html": 1, "md": 0, "json": 0, "csv": 0, "txt": 2},
    "04-sales-pre": {"docx": 3, "pdf": 2, "pptx": 5, "xlsx": 1, "html": 1, "md": 0, "json": 0, "csv": 0, "txt": 1},
    "05-faq-support": {"docx": 2, "pdf": 0, "pptx": 0, "xlsx": 4, "html": 3, "md": 0, "json": 3, "csv": 3, "txt": 2},
    "06-project-delivery": {"docx": 5, "pdf": 3, "pptx": 0, "xlsx": 3, "html": 1, "md": 0, "json": 0, "csv": 1, "txt": 2},
    "07-training-hr": {"docx": 3, "pdf": 0, "pptx": 5, "xlsx": 0, "html": 2, "md": 0, "json": 0, "csv": 0, "txt": 2},
    "08-meeting-decision": {"docx": 5, "pdf": 0, "pptx": 0, "xlsx": 2, "html": 2, "md": 0, "json": 0, "csv": 0, "txt": 2},
    "09-version-conflict": {"docx": 4, "pdf": 3, "pptx": 0, "xlsx": 3, "html": 1, "md": 0, "json": 0, "csv": 0, "txt": 2},
    "10-noise": {"txt": 5},
}

ALL_FORMATS = ["docx", "pdf", "pptx", "xlsx", "html", "txt", "md", "json", "csv"]


def _is_heading(line: str) -> bool:
    s = line.strip()
    if not s:
        return False
    if s.startswith(("【", "#")):
        return True
    if s.startswith(("一、", "二、", "三、", "四、", "五、", "六、", "七、", "八、", "九、", "十、")):
        return True
    if s.startswith("第") and "条" in s[:8]:
        return True
    if s.startswith(("第一章", "第二章", "第三章", "第四章", "第五章", "第六章")):
        return True
    return False


def _write_docx(text: str, path: Path) -> None:
    from docx import Document

    doc = Document()
    for line in text.splitlines():
        s = line.strip()
        if not s:
            continue
        if _is_heading(s):
            doc.add_heading(s, level=2)
        else:
            doc.add_paragraph(s)
    doc.save(str(path))


def _write_pdf(text: str, path: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    style = ParagraphStyle("cn", fontName="STSong-Light", fontSize=11, leading=16)
    title = ParagraphStyle("cn-title", fontName="STSong-Light", fontSize=16, leading=22, spaceAfter=10)
    doc = SimpleDocTemplate(str(path), pagesize=A4)
    story: list[object] = []
    first = True
    for line in text.splitlines():
        s = line.strip()
        if not s:
            story.append(Spacer(1, 6))
            continue
        escaped = html.escape(s)
        if first:
            story.append(Paragraph(escaped, title))
            first = False
        elif _is_heading(s):
            story.append(Paragraph(escaped, ParagraphStyle("h", fontName="STSong-Light", fontSize=13, leading=18, spaceBefore=8)))
        else:
            story.append(Paragraph(escaped, style))
    doc.build(story)


def _write_pptx(text: str, path: Path) -> None:
    from pptx import Presentation

    lines = [l.strip() for l in text.splitlines() if l.strip()]
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = (lines[0][:40] if lines else "文档")
    chunk: list[str] = []

    def flush() -> None:
        if not chunk:
            return
        s = prs.slides.add_slide(prs.slide_layouts[1])
        body = s.shapes.placeholders[1]
        tf = body.text_frame
        first = True
        for line in chunk:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = line[:90]
            first = False
        chunk.clear()

    for line in lines[1:]:
        chunk.append(line)
        if len(chunk) >= 8:
            flush()
    flush()
    prs.save(str(path))


def _first_line(text: str) -> str:
    for line in text.splitlines():
        s = line.strip()
        if s:
            return s
    return "文档"


def _parse_doc(text: str) -> dict:
    """Parse a generated doc into structured parts for tabular formats.

    Returns {title, metadata:[[k,v],...], qa:[[q,a],...], tables:[row,...],
    sections:[[heading, body],...]}.
    """
    metadata: list[list[str]] = []
    qa: list[list[str]] = []
    tables: list[list[str]] = []
    sections: list[list[str]] = []
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    title = lines[0] if lines else ""
    header_mode = True
    cur_heading: str | None = None
    cur_body: list[str] = []
    pending_q: list[str] | None = None

    for s in lines[1:]:
        # header metadata block: key：value before the first section heading
        if header_mode and ("：" in s) and not _is_heading(s):
            k, v = s.split("：", 1)
            if len(k) <= 10 and not s.startswith(("Q", "A")):
                metadata.append([k.strip(), v.strip()])
                continue
        if _is_heading(s):
            header_mode = False
            if pending_q:
                qa.append(pending_q)
                pending_q = None
            if cur_heading:
                sections.append([cur_heading, "\n".join(cur_body)])
            cur_heading, cur_body = s, []
            continue
        header_mode = False
        # FAQ Q / A
        if re.match(r"^Q\d?：", s):
            if pending_q:
                qa.append(pending_q)
            pending_q = [s.split("：", 1)[1].strip(), ""]
            continue
        if pending_q is not None and (s.startswith("A") or s.startswith("答")):
            pending_q[1] = s.split("：", 1)[1].strip() if "：" in s else s
            continue
        if pending_q is not None:
            # 遇到"补充说明"等结尾标记时结束问答块
            if s.startswith("补充") or s.startswith("（本文档"):
                qa.append(pending_q)
                pending_q = None
                cur_body.append(s)
                continue
            pending_q[1] = (pending_q[1] + " " + s).strip()
            continue
        # markdown table rows
        if s.startswith("|") or s.count("|") >= 2:
            cells = [c.strip() for c in s.strip("|").split("|")]
            while cells and not cells[0]:
                cells.pop(0)
            while cells and not cells[-1]:
                cells.pop()
            if cells and not all(c == "-" or c == "---" for c in cells):
                tables.append(cells)
            continue
        cur_body.append(s)

    if pending_q:
        qa.append(pending_q)
    if cur_heading:
        sections.append([cur_heading, "\n".join(cur_body)])
    return {"title": title, "metadata": metadata, "qa": qa, "tables": tables, "sections": sections}


def _write_xlsx(text: str, path: Path) -> None:
    from openpyxl import Workbook

    doc = _parse_doc(text)
    wb = Workbook()
    ws = wb.active
    ws.title = "元数据"
    ws.append(["字段", "内容"])
    ws.append(["标题", doc["title"]])
    for k, v in doc["metadata"]:
        ws.append([k, v])
    if doc["qa"]:
        ws2 = wb.create_sheet("问答")
        ws2.append(["问题", "答案"])
        for q, a in doc["qa"]:
            ws2.append([q, a])
    if doc["tables"]:
        ws3 = wb.create_sheet("表格")
        for row in doc["tables"]:
            ws3.append(row)
    if doc["sections"]:
        ws4 = wb.create_sheet("内容")
        ws4.append(["章节", "内容"])
        for h, b in doc["sections"]:
            ws4.append([h, b[:500]])
    wb.save(str(path))


def _write_csv(text: str, path: Path) -> None:
    import csv as _csv

    doc = _parse_doc(text)
    with open(path, "w", newline="", encoding="utf-8-sig") as fp:
        w = _csv.writer(fp)
        # 永远先输出含唯一“文档编号”的元数据表，保证 CSV 不重复
        w.writerow(["字段", "内容"])
        w.writerow(["标题", doc["title"]])
        for k, v in doc["metadata"]:
            w.writerow([k, v])
        if doc["qa"]:
            w.writerow(["问题", "答案"])
            for q, a in doc["qa"]:
                w.writerow([q, a])
        for row in doc["tables"]:
            w.writerow(row)
        if doc["sections"]:
            w.writerow(["章节", "内容"])
            for h, b in doc["sections"]:
                w.writerow([h, b[:300]])


def _write_json(text: str, path: Path) -> None:
    doc = _parse_doc(text)
    data = {
        "title": doc["title"],
        "metadata": dict(doc["metadata"]),
        "qa": [{"question": q, "answer": a} for q, a in doc["qa"]],
        "content": text,
    }
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _write_html(text: str, path: Path) -> None:
    body_lines = []
    for line in text.splitlines():
        s = line.strip()
        if not s:
            continue
        if _is_heading(s):
            body_lines.append(f"<h3>{html.escape(s)}</h3>")
        elif s.startswith("|") or "| " in s:
            cells = [c.strip() for c in s.strip("|").split("|")]
            body_lines.append("<tr>" + "".join(f"<td>{html.escape(c)}</td>" for c in cells) + "</tr>")
        else:
            body_lines.append(f"<p>{html.escape(s)}</p>")
    title = html.escape(_first_line(text)[:60])
    doc = (
        "<!DOCTYPE html><html lang=\"zh-CN\"><head><meta charset=\"utf-8\">"
        f"<title>{title}</title></head><body>"
        + "".join(body_lines)
        + "</body></html>"
    )
    path.write_text(doc, encoding="utf-8")


def _write_md(text: str, path: Path) -> None:
    # Content already uses markdown-style headings; emit as-is with a title.
    path.write_text(text, encoding="utf-8")


WRITERS = {"docx": _write_docx, "pdf": _write_pdf, "pptx": _write_pptx, "xlsx": _write_xlsx,
           "html": _write_html, "md": _write_md, "json": _write_json, "csv": _write_csv, "txt": _write_md}


def render_file(text: str, fmt: str, path: Path) -> None:
    """Render plain text content into a file of the given format."""
    WRITERS[fmt](text, path)
